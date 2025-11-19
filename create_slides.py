from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # --- Helper Functions ---
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = subtitle

    def add_section_header(title):
        slide_layout = prs.slide_layouts[2] # Section Header
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

    def add_content_slide(title, content_bullets, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for bullet in content_bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.space_after = Pt(10)
            p.level = 0
            # Handle sub-bullets (simple heuristic: starts with "  ")
            if bullet.startswith("  "):
                p.text = bullet.strip()
                p.level = 1
            elif bullet.startswith("    "):
                p.text = bullet.strip()
                p.level = 2

        # Image
        if image_path and os.path.exists(image_path):
            # Resize text box
            body_shape.width = Inches(5.5)
            
            # Add image
            left = Inches(5.8)
            top = Inches(2.0)
            height = Inches(3.5)
            slide.shapes.add_picture(image_path, left, top, height=height)

    # --- Image Paths ---
    artifacts_dir = r"C:\Users\加藤裕樹\.gemini\antigravity\brain\26ba0119-23a2-42e1-9a6d-dd9db635e993"
    img_nmf = os.path.join(artifacts_dir, "nmf_concept_1763548181387.png")
    img_graph = os.path.join(artifacts_dir, "graph_concept_1763548189193.png")
    img_sem = os.path.join(artifacts_dir, "sem_concept_1763548205134.png")
    img_lingam = os.path.join(artifacts_dir, "lingam_concept_1763548214596.png")

    # --- Main Presentation ---
    
    add_title_slide(
        "推薦アルゴリズムの数理と実装",
        "CareerNavigatorにおける4つのアプローチ\n(NMF, Graph, SEM, LiNGAM)"
    )

    # ==========================================
    # 1. NMF (Non-negative Matrix Factorization)
    # ==========================================
    add_section_header("1. NMF (非負値行列因子分解)")

    # 1-1. Overview
    add_content_slide(
        "1.1 NMF: 概要と目的",
        [
            "目的: 高次元の疎なユーザー・アイテム行列から、潜在的な特徴を抽出する",
            "アプローチ: 行列分解 (Matrix Factorization) の一種",
            "特徴: すべての要素が非負 (≧0) であるという制約",
            "解釈性: 「部分」の加算として全体を表現するため、PCA等と比較して解釈が容易",
            "応用: 協調フィルタリング、トピックモデル、画像の特徴抽出"
        ]
    )
    # 1-2. Problem Formulation
    add_content_slide(
        "1.2 NMF: 問題定式化",
        [
            "入力: 非負行列 V (m × n)",
            "  - m: ユーザー数 (メンバー)",
            "  - n: アイテム数 (スキル)",
            "  - V_ij: ユーザ i のスキル j に対するスコア (習得レベル等)",
            "出力: 2つの非負行列 W, H",
            "  - W (m × r): 基底行列 (メンバー特性)",
            "  - H (r × n): 係数行列 (スキル特性)",
            "  - r: ランク (潜在因子数, r << min(m, n))",
            "近似式: V ≈ WH"
        ]
    )
    # 1-3. Mathematical Model
    add_content_slide(
        "1.3 NMF: 数理モデル",
        [
            "以下の目的関数を最小化する最適化問題として定式化される:",
            "Objective Function (Frobenius Norm):",
            "  min || V - WH ||_F^2",
            "  subject to W, H ≧ 0",
            "ここで、フロベニウスノルムは:",
            "  || A ||_F = √(Σ Σ |a_ij|^2)",
            "別解: Kullback-Leibler Divergence (一般化KL情報量) を用いる場合もある"
        ]
    )
    # 1-4. Optimization Algorithm (1)
    add_content_slide(
        "1.4 NMF: 最適化アルゴリズム (乗法更新則)",
        [
            "非凸最適化問題であるため、大域的最適解は保証されない",
            "一般的に Lee & Seung (2001) の Multiplicative Update Rules を使用",
            "勾配降下法と異なり、学習率パラメータが不要で、非負制約を自然に満たす",
            "更新式:",
            "  H ← H .* (W^T V) ./ (W^T W H)",
            "  W ← W .* (V H^T) ./ (W H H^T)",
            "  (.* は要素ごとの積、./ は要素ごとの商)"
        ]
    )
    # 1-5. Optimization Algorithm (2)
    add_content_slide(
        "1.5 NMF: 収束と停止条件",
        [
            "反復計算の流れ:",
            "  1. W, H を非負の乱数で初期化",
            "  2. H を更新",
            "  3. W を更新",
            "  4. 目的関数の変化が閾値以下になるまで 2-3 を繰り返す",
            "局所解 (Local Minima) に陥る可能性があるため、初期値依存性がある",
            "対策: 複数回の初期化で試行する、SVDベースの初期化 (NNDSVD) を用いる"
        ]
    )
    # 1-6. Key Features
    add_content_slide(
        "1.6 NMF: アルゴリズムの特性",
        [
            "スパース性 (Sparsity):",
            "  - 非負制約により、W, H の多くの要素が0になりやすい",
            "  - 結果として、解釈可能な「部品」が抽出される",
            "部分による全体表現 (Parts-based Representation):",
            "  - 負の値（打ち消し合い）がないため、構成要素の「足し合わせ」のみで表現",
            "  - 例: 「顔」画像 = 「目」+「鼻」+「口」の基底画像の和"
        ]
    )
    # 1-7. Implementation Details
    add_content_slide(
        "1.7 NMF: 本システムでの実装",
        [
            "ライブラリ: scikit-learn (sklearn.decomposition.NMF)",
            "パラメータ設定:",
            "  - n_components (ランク): 20 (実験的に決定)",
            "  - init: 'nndsvda' (NNDSVD with zeros filled with average)",
            "  - solver: 'cd' (Coordinate Descent) または 'mu' (Multiplicative Update)",
            "前処理:",
            "  - 欠損値の補完 (0埋め)",
            "  - 正規化 (TF-IDF的な重み付けも検討可能)"
        ]
    )
    # 1-8. Recommendation Logic
    add_content_slide(
        "1.8 NMF: 推薦ロジック",
        [
            "スコア予測:",
            "  Predicted_V = W × H",
            "  ユーザー u の未習得スキル i に対するスコア = (W_u) ・ (H_i)^T",
            "推薦リスト生成:",
            "  1. 再構成行列を計算",
            "  2. 既習得スキルを除外（マスク）",
            "  3. スコアが高い順にトップNを抽出",
            "解釈: 「あなたの潜在特性ベクトル W_u は、このスキル H_i と相性が良い」"
        ]
    )
    # 1-9. Visual Concept
    add_content_slide(
        "1.9 NMF: 概念図",
        [
            "巨大な行列を2つの小さな行列に分解するイメージ",
            "左側: ユーザーの潜在特性",
            "右側: スキルの潜在特性",
            "その積が元のデータを近似する"
        ],
        img_nmf
    )
    # 1-10. Pros & Cons
    add_content_slide(
        "1.10 NMF: 利点と課題",
        [
            "利点:",
            "  - データ駆動で未知のパターンを発見できる (Serendipity)",
            "  - 類似ユーザー・類似アイテムの情報を同時に活用できる",
            "  - 実装が比較的容易で高速",
            "課題:",
            "  - コールドスタート問題 (新規ユーザー・スキルへの対応困難)",
            "  - 説明可能性が限定的 (「なぜこの潜在因子なのか」の説明が難しい)",
            "  - スパース性が極端に高いデータでは精度が低下する"
        ]
    )

    # ==========================================
    # 2. Graph (Skill Transition Graph)
    # ==========================================
    add_section_header("2. Graph (スキル遷移グラフ & Node2Vec)")

    # 2-1. Overview
    add_content_slide(
        "2.1 Graph: 概要と目的",
        [
            "目的: スキルの習得順序（遷移）をグラフ構造としてモデル化し、最適な学習パスを提案する",
            "アプローチ: グラフ理論 + グラフ埋め込み (Graph Embedding)",
            "特徴: 「順序性」と「近接性」を明示的に扱う",
            "応用: 経路探索、次アイテム推薦、リンク予測"
        ]
    )
    # 2-2. Graph Construction
    add_content_slide(
        "2.2 Graph: グラフ構築",
        [
            "ノード (V): スキル",
            "エッジ (E): 習得の遷移 (スキルA → スキルB)",
            "エッジ重み (Weight):",
            "  - 遷移人数 (Frequency)",
            "  - 時間減衰 (Time Decay): 最近の遷移を重視",
            "  - w_ij = Σ exp(-λ * Δt)",
            "有向グラフ (Directed Graph) として構築"
        ]
    )
    # 2-3. Node2Vec Algorithm (1)
    add_content_slide(
        "2.3 Node2Vec: アルゴリズム概要",
        [
            "目的: グラフ上のノードを低次元ベクトル空間に写像する",
            "基本思想: Skip-gram (Word2Vec) のグラフへの拡張",
            "  - 「グラフ上で近くにあるノードは、ベクトル空間でも近くにあるべき」",
            "手順:",
            "  1. ランダムウォークでノード列（文脈）を生成",
            "  2. Skip-gramモデルでベクトル表現を学習"
        ]
    )
    # 2-4. Node2Vec Algorithm (2) - Random Walk
    add_content_slide(
        "2.4 Node2Vec: Biased Random Walk",
        [
            "探索の制御パラメータ p, q:",
            "  - Return parameter (p): 前のノードに戻る確率 (BFS的挙動、局所構造)",
            "  - In-out parameter (q): 遠くへ行く確率 (DFS的挙動、大域構造)",
            "柔軟性:",
            "  - p, q の調整により、同質性 (Homophily) と 構造的等価性 (Structural Equivalence) のバランスをとる",
            "本システム設定: q > 1 (DFS寄り) でパスの連続性を重視"
        ]
    )
    # 2-5. Mathematical Model (Skip-gram)
    add_content_slide(
        "2.5 Node2Vec: 数理モデル (Skip-gram)",
        [
            "目的関数 (対数尤度の最大化):",
            "  max Σ_{u ∈ V} Σ_{v ∈ N_R(u)} log P(v | u)",
            "  - N_R(u): ランダムウォークで生成された u の近傍ノード集合",
            "Softmax確率:",
            "  P(v | u) = exp(u・v) / Σ_{n ∈ V} exp(u・n)",
            "計算効率化:",
            "  - Negative Sampling を用いて分母の計算コストを削減"
        ]
    )
    # 2-6. Recommendation Logic (1) - Graph Traversal
    add_content_slide(
        "2.6 Graph: 推薦ロジック (1) - 直接遷移",
        [
            "アプローチ: グラフ上の直接的なつながりを評価",
            "スコア計算:",
            "  Score(Target) = Σ_{Owned} Weight(Owned → Target)",
            "特徴:",
            "  - 実績に基づく確実な「次の一歩」",
            "  - エッジが存在しない（誰も遷移していない）ペアは推薦されない"
        ]
    )
    # 2-7. Recommendation Logic (2) - Embedding Similarity
    add_content_slide(
        "2.7 Graph: 推薦ロジック (2) - 埋め込み類似度",
        [
            "アプローチ: Node2Vecで得られたベクトル間の類似度",
            "指標: コサイン類似度 (Cosine Similarity)",
            "  Sim(A, B) = (A・B) / (||A|| ||B||)",
            "特徴:",
            "  - 直接の遷移がなくても、文脈（周辺構造）が似ていれば推薦可能",
            "  - 潜在的な関連性を捕捉"
        ]
    )
    # 2-8. Implementation Details
    add_content_slide(
        "2.8 Graph: 本システムでの実装",
        [
            "ライブラリ: networkx, node2vec",
            "パラメータ:",
            "  - dimensions: 64 (ベクトル次元数)",
            "  - walk_length: 10 (ウォーク長)",
            "  - num_walks: 80 (各ノードからの試行回数)",
            "  - window: 5 (Skip-gramのウィンドウサイズ)",
            "ハイブリッド推薦:",
            "  直接遷移スコアと埋め込み類似度スコアを重み付け統合"
        ]
    )
    # 2-9. Visual Concept
    add_content_slide(
        "2.9 Graph: 概念図",
        [
            "ノード（スキル）とエッジ（遷移）によるネットワーク構造",
            "太いエッジは多くの人が通った「王道ルート」",
            "ランダムウォークにより、グラフ構造をベクトル化する"
        ],
        img_graph
    )
    # 2-10. Pros & Cons
    add_content_slide(
        "2.10 Graph: 利点と課題",
        [
            "利点:",
            "  - 「順序」を考慮できる（NMFにはない特徴）",
            "  - 説明可能性が高い（「Aを学んだ人はBを学んでいます」）",
            "  - ネットワーク分析の指標（中心性など）も活用可能",
            "課題:",
            "  - データが少ないとグラフが分断される",
            "  - 過去のトレンドに引きずられやすい（ポピュラリティバイアス）",
            "  - 計算コスト（特に大規模グラフでの埋め込み学習）"
        ]
    )

    # ==========================================
    # 3. SEM (Structural Equation Modeling)
    # ==========================================
    add_section_header("3. SEM (構造方程式モデリング)")

    # 3-1. Overview
    add_content_slide(
        "3.1 SEM: 概要と目的",
        [
            "目的: 観測変数（スキル）の背後にある潜在変数（能力・因子）と、その因果構造をモデル化する",
            "アプローチ: 共分散構造分析",
            "特徴: 仮説検証型のアプローチ（モデル構造を人間が定義）",
            "応用: 心理統計、社会科学、マーケティング、能力評価"
        ]
    )
    # 3-2. Model Structure
    add_content_slide(
        "3.2 SEM: モデルの構成要素",
        [
            "1. 測定モデル (Measurement Model):",
            "  - 潜在変数が観測変数をどう規定するか",
            "  - 因子分析に相当",
            "  - 例: 「基礎力」→「Python」「SQL」",
            "2. 構造モデル (Structural Model):",
            "  - 潜在変数間の因果関係",
            "  - 回帰分析の拡張",
            "  - 例: 「基礎力」→「応用力」"
        ]
    )
    # 3-3. Mathematical Formulation
    add_content_slide(
        "3.3 SEM: 数理モデル",
        [
            "構造方程式:",
            "  η = Bη + Γξ + ζ",
            "  - η: 内生潜在変数, ξ: 外生潜在変数",
            "  - B, Γ: パス係数行列, ζ: 誤差項",
            "測定方程式:",
            "  y = Λ_y η + ε",
            "  x = Λ_x ξ + δ",
            "  - x, y: 観測変数",
            "  - Λ: 因子負荷行列, ε, δ: 測定誤差"
        ]
    )
    # 3-4. Parameter Estimation
    add_content_slide(
        "3.4 SEM: パラメータ推定",
        [
            "目的: モデルから導かれる共分散行列 Σ(θ) が、データの標本共分散行列 S に近づくようにパラメータ θ を推定",
            "推定法: 最尤推定法 (Maximum Likelihood Estimation: ML)",
            "適合関数 (Fitting Function):",
            "  F_ML = log|Σ(θ)| + tr(S Σ(θ)^-1) - log|S| - p",
            "  - p: 変数の数",
            "  - 多変量正規分布を仮定"
        ]
    )
    # 3-5. Model Evaluation
    add_content_slide(
        "3.5 SEM: モデル適合度指標",
        [
            "モデルの良さを評価する指標:",
            "GFI (Goodness of Fit Index):",
            "  - 決定係数のようなもの。0.9以上が目安。",
            "RMSEA (Root Mean Square Error of Approximation):",
            "  - 近似誤差の平均。0.05以下が良い。",
            "CFI (Comparative Fit Index):",
            "  - 独立モデルとの比較。0.95以上が良い。",
            "AIC (Akaike Information Criterion):",
            "  - モデル比較用（小さいほど良い）"
        ]
    )
    # 3-6. Implementation Details
    add_content_slide(
        "3.6 SEM: 本システムでの実装",
        [
            "実装: Pythonでのスクラッチ実装 (UnifiedSEMEstimator)",
            "最適化: scipy.optimize.minimize を使用",
            "  - 制約付き最適化 (分散≧0 等)",
            "モデル定義:",
            "  - キャリアパス階層に基づき、測定モデルと構造モデルを動的に生成",
            "  - 例: Stage1(潜在) → Stage2(潜在) というパスを固定"
        ]
    )
    # 3-7. Recommendation Logic
    add_content_slide(
        "3.7 SEM: 推薦ロジック",
        [
            "因子スコアの推定:",
            "  - メンバーの観測スキルから、潜在変数（能力値）を推定",
            "ギャップ分析:",
            "  - 目標とするキャリアステージ（潜在変数）に必要なスコアと現状の差分",
            "推薦:",
            "  - 因子負荷量 (Λ) が高く、かつ未習得のスキルを優先",
            "  - 「あなたの基礎力向上には、このスキルが最も寄与します」"
        ]
    )
    # 3-8. Advantages
    add_content_slide(
        "3.8 SEM: 利点",
        [
            "理論的妥当性:",
            "  - 専門知識（ドメイン知識）をモデル構造に反映できる",
            "  - 「なぜ」の構造が明確",
            "潜在変数の扱い:",
            "  - 直接測定できない「能力」や「概念」を扱える",
            "  - 測定誤差を明示的にモデル化できるため、パラメータ推定のバイアスが小さい"
        ]
    )
    # 3-9. Visual Concept
    add_content_slide(
        "3.9 SEM: 概念図",
        [
            "潜在変数（丸）と観測変数（四角）のパス図",
            "矢印は因果（回帰）関係を表す",
            "階層的な構造を持つことが多い"
        ],
        img_sem
    )
    # 3-10. Limitations
    add_content_slide(
        "3.10 SEM: 課題と限界",
        [
            "モデル特定化の困難さ:",
            "  - 適切なモデル構造を事前に定義する必要がある",
            "  - 探索的な利用には向かない（修正指数などを使うが限界あり）",
            "計算コスト:",
            "  - 変数が増えると計算が不安定になりやすい",
            "  - 収束しない場合がある",
            "分布の仮定:",
            "  - 基本的に多変量正規分布を仮定するため、離散データ（有無）への適用には工夫が必要"
        ]
    )

    # ==========================================
    # 4. LiNGAM (Linear Non-Gaussian Acyclic Model)
    # ==========================================
    add_section_header("4. LiNGAM (線形非ガウス非巡回モデル)")

    # 4-1. Overview
    add_content_slide(
        "4.1 LiNGAM: 概要と目的",
        [
            "目的: 観測データのみから、変数間の因果構造（向き）を識別する",
            "アプローチ: 因果探索 (Causal Discovery)",
            "特徴: 従来の統計手法（相関）では不可能な「因果の方向」の特定が可能",
            "  - 「相関関係」≠「因果関係」の壁を突破",
            "応用: 経済分析、生物学（遺伝子ネットワーク）、マーケティング"
        ]
    )
    # 4-2. The Identification Problem
    add_content_slide(
        "4.2 LiNGAM: 識別性の問題",
        [
            "なぜ因果探索は難しいか？",
            "  - ガウス分布（正規分布）の世界では、X→Y と Y→X の区別がつかない",
            "  - どちらのモデルでも同じ共分散行列を生成できてしまう（モデルの等価性）",
            "LiNGAMのブレイクスルー:",
            "  - 「データが非ガウス分布であれば、因果の方向を一意に特定できる」ことを証明 (Shimizu et al., 2006)"
        ]
    )
    # 4-3. Assumptions & Model
    add_content_slide(
        "4.3 LiNGAM: 仮定と数理モデル",
        [
            "3つの仮定:",
            "  1. 線形性 (Linearity): x_i = Σ b_ij x_j + e_i",
            "  2. 非ガウス性 (Non-Gaussianity): 誤差項 e_i は非正規分布に従う",
            "  3. 非巡回性 (Acyclicity): 因果グラフはDAG (有向非巡回グラフ) である",
            "行列表記:",
            "  x = Bx + e",
            "  x = (I - B)^-1 e = A e  (A: 混合行列)"
        ]
    )
    # 4-4. Algorithm (ICA-based)
    add_content_slide(
        "4.4 LiNGAM: アルゴリズム (ICAベース)",
        [
            "基本手順:",
            "  1. 独立成分分析 (ICA) を行い、混合行列 A を推定",
            "     x = A s  (sは独立成分)",
            "  2. 行列 A の行・列を並べ替えて、対角成分が非ゼロになるように正規化",
            "  3. さらに並べ替えて、下三角行列に近づける",
            "     (因果的順序の決定: 原因 → 結果)",
            "  4. 行列 B を計算し、有意でない係数を剪定 (Pruning)"
        ]
    )
    # 4-5. DirectLiNGAM
    add_content_slide(
        "4.5 DirectLiNGAM: 改良アルゴリズム",
        [
            "ICAベースの課題: 初期値依存性、収束性",
            "DirectLiNGAM (2011):",
            "  - ICAを使わず、変数間の独立性を直接評価して因果順序を決定",
            "手順:",
            "  1. すべての変数に対して、他の変数との独立性を評価",
            "  2. 最も「外生的な（原因となる）」変数を特定し、順序リストに追加",
            "  3. その変数の影響（成分）を他の変数から除去（回帰残差）",
            "  4. 残りの変数で 1-3 を繰り返す",
            "特徴: 決定論的で必ず解が得られる"
        ]
    )
    # 4-6. Implementation Details
    add_content_slide(
        "4.6 LiNGAM: 本システムでの実装",
        [
            "ライブラリ: lingam (Python package)",
            "使用モデル: DirectLiNGAM",
            "計算コスト対策:",
            "  - 変数（スキル）が多すぎると計算量が爆発する (O(p^3) ~ O(p^4))",
            "  - クラスタリングによる分割統治法を採用",
            "  - 類似スキル群ごとにLiNGAMを適用し、全体構造を統合"
        ]
    )
    # 4-7. Recommendation Logic
    add_content_slide(
        "4.7 LiNGAM: 推薦ロジック",
        [
            "因果効果 (Causal Effect) の算出:",
            "  - Total Effect = 直接効果 + 間接効果",
            "  - 行列 B から (I - B)^-1 を計算することで全効果を取得",
            "スコアリング:",
            "  1. Readiness Score: 既習得スキルが、ターゲットスキルの習得をどれだけ促進するか（原因→結果）",
            "  2. Utility Score: ターゲットスキルが、将来のスキル習得にどれだけ役立つか（結果→原因）",
            "  - これらを統合してランク付け"
        ]
    )
    # 4-8. Visual Concept
    add_content_slide(
        "4.8 LiNGAM: 概念図",
        [
            "虫眼鏡で「真の矢印」を見つけるイメージ",
            "相関関係（双方向）ではなく、因果関係（片方向）を特定",
            "データの分布形状（非ガウス性）が鍵"
        ],
        img_lingam
    )
    # 4-9. Pros
    add_content_slide(
        "4.9 LiNGAM: 利点",
        [
            "因果の特定:",
            "  - 「AをやればBができるようになる」という介入効果を予測できる",
            "  - 推薦の説得力が高い（根拠が明確）",
            "モデルフリー:",
            "  - SEMのように事前に構造を指定する必要がない",
            "  - データから構造自体を発見できる"
        ]
    )
    # 4-10. Cons & Challenges
    add_content_slide(
        "4.10 LiNGAM: 課題と限界",
        [
            "仮定の制約:",
            "  - 非ガウス性が必要（データが正規分布に近いと機能しない）",
            "  - 非巡回（フィードバックループを持てない）",
            "  - 未観測共通原因（交絡因子）の影響を受けやすい",
            "データ量と計算量:",
            "  - 信頼性の高い推定には多くのサンプルが必要",
            "  - 変数が多いと計算時間が長い"
        ]
    )

    # ==========================================
    # Summary
    # ==========================================
    add_section_header("まとめ")
    
    add_content_slide(
        "アルゴリズムの比較と使い分け",
        [
            "NMF (協調フィルタリング):",
            "  - 探索的発見重視。データがスパースでも機能しやすい。",
            "Graph (遷移グラフ):",
            "  - 実績重視。学習パスの順序性を自然に表現。",
            "SEM (構造方程式):",
            "  - 理論重視。ドメイン知識を組み込み、仮説を検証する。",
            "LiNGAM (因果探索):",
            "  - 因果重視。データから真の前提条件構造を発見する。",
            "結論: 目的に応じてこれらをハイブリッドに組み合わせることが有効"
        ]
    )

    prs.save('Recommendation_Algorithms_Explanation.pptx')
    print("Presentation saved as Recommendation_Algorithms_Explanation.pptx")

if __name__ == "__main__":
    create_presentation()
